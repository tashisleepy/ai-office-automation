#!/usr/bin/env python3
"""
python-pptx Template Generator
Design a master template once → auto-fill with any data forever.
Feed it client data, forensic analysis, or AI output — get consulting-grade decks.

Usage:
  python3 scripts/pptx-template.py
  python3 scripts/pptx-template.py --output my-deck.pptx

Dependencies: pip3 install python-pptx
"""

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── BRAND CONFIG ──
BRAND = {
    "primary": RGBColor(0x1E, 0x27, 0x61),     # Navy
    "secondary": RGBColor(0xCA, 0xDC, 0xFC),    # Ice blue
    "accent": RGBColor(0xFF, 0xFF, 0xFF),        # White
    "bg_light": RGBColor(0xF5, 0xF5, 0xF5),     # Light gray
    "text_dark": RGBColor(0x33, 0x33, 0x33),     # Dark gray
    "text_muted": RGBColor(0x99, 0x99, 0x99),    # Muted
    "font_title": "Arial",
    "font_body": "Calibri",
}


def add_title_slide(prs, title, subtitle):
    """Dark navy title slide with accent line."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BRAND["primary"]

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = BRAND["accent"]
    p.font.name = BRAND["font_title"]

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.italic = True
    p2.font.color.rgb = BRAND["secondary"]
    p2.font.name = BRAND["font_body"]

    # Accent line
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.2), Inches(2.5), Pt(3)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BRAND["secondary"]
    slide.shapes[-1].line.fill.background()

    return slide


def add_content_slide(prs, heading, items):
    """Light background slide with numbered items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BRAND["bg_light"]

    # Heading
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = heading.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BRAND["primary"]
    p.font.name = BRAND["font_title"]

    # Items with numbered circles
    for i, item in enumerate(items):
        y_pos = 2.0 + (i * 1.5)

        # Number circle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), Inches(y_pos), Inches(0.5), Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BRAND["primary"]
        shape.line.fill.background()
        tf_num = shape.text_frame
        tf_num.paragraphs[0].text = str(i + 1)
        tf_num.paragraphs[0].font.size = Pt(16)
        tf_num.paragraphs[0].font.bold = True
        tf_num.paragraphs[0].font.color.rgb = BRAND["accent"]
        tf_num.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_num.word_wrap = False

        # Title
        txBox_t = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos - 0.05), Inches(10), Inches(0.4))
        p_t = txBox_t.text_frame.paragraphs[0]
        p_t.text = item["title"]
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = BRAND["primary"]
        p_t.font.name = BRAND["font_body"]

        # Description
        txBox_d = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos + 0.35), Inches(10), Inches(0.4))
        p_d = txBox_d.text_frame.paragraphs[0]
        p_d.text = item["desc"]
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = BRAND["text_dark"]
        p_d.font.name = BRAND["font_body"]

    return slide


def add_table_slide(prs, heading, headers, rows):
    """Slide with a styled data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BRAND["bg_light"]

    # Heading
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = heading.upper()
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRAND["primary"]
    p.font.name = BRAND["font_title"]

    # Table
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0)
    )
    table = table_shape.table

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND["primary"]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = BRAND["accent"]
            paragraph.font.name = BRAND["font_body"]

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = BRAND["text_dark"]
                paragraph.font.name = BRAND["font_body"]

    return slide


def add_contact_slide(prs, email, phone=None, location=None):
    """Dark navy contact slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BRAND["primary"]

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0))
    p = txBox.text_frame.paragraphs[0]
    p.text = "GET IN TOUCH"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = BRAND["accent"]
    p.font.name = BRAND["font_title"]
    p.alignment = PP_ALIGN.CENTER

    y = 3.0
    for text in [email, phone, location]:
        if text:
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.6))
            p2 = txBox2.text_frame.paragraphs[0]
            p2.text = text
            p2.font.size = Pt(22)
            p2.font.color.rgb = BRAND["secondary"]
            p2.font.name = BRAND["font_body"]
            p2.alignment = PP_ALIGN.CENTER
            y += 0.7

    return slide


def build_deck(data, output="output-deck.pptx"):
    """Build a full deck from a data dictionary."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    add_title_slide(prs, data["title"], data["subtitle"])

    # Content slides
    for slide_data in data.get("slides", []):
        if "items" in slide_data:
            add_content_slide(prs, slide_data["heading"], slide_data["items"])
        elif "table" in slide_data:
            add_table_slide(prs, slide_data["heading"],
                          slide_data["table"]["headers"],
                          slide_data["table"]["rows"])

    # Contact
    if "contact" in data:
        add_contact_slide(prs, **data["contact"])

    prs.save(output)
    print(f"Deck created: {output}")


# ── EXAMPLE DATA ──
if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", "-o", default="output-deck.pptx")
    args = parser.parse_args()

    example_data = {
        "title": "Your Company",
        "subtitle": "Enterprise AI Solutions",
        "slides": [
            {
                "heading": "Services",
                "items": [
                    {"title": "AI Consulting", "desc": "End-to-end AI strategy and deployment for enterprise"},
                    {"title": "Competitive Intelligence", "desc": "Forensic teardowns with funding reality checks"},
                    {"title": "Product Development", "desc": "Production-ready AI systems deployed at scale"},
                ]
            },
            {
                "heading": "Client Portfolio",
                "table": {
                    "headers": ["Client", "Industry", "Engagement", "Value"],
                    "rows": [
                        ["Media Corp A", "Media", "AI Video Pipeline", "Enterprise"],
                        ["Media Corp B", "Media", "Content Automation", "Enterprise"],
                        ["Sports Franchise", "Sports", "AI Operations", "Enterprise"],
                        ["Retail Group", "Retail", "AI Systems", "Enterprise"],
                        ["Energy Corp", "Energy", "AI Solutions", "Enterprise"],
                    ]
                }
            }
        ],
        "contact": {
            "email": "hello@example.com",
            "phone": "+91 XXXXXXXXXX",
            "location": "Your City, Your Country"
        }
    }

    build_deck(example_data, args.output)
